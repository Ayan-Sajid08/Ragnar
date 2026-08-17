from io import BytesIO

from pptx import Presentation

from services.chunker import chunk_text


def extract_text(file_bytes: bytes):
    presentation = Presentation(BytesIO(file_bytes))

    text_chunks = []
    chunk_index = 0

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_text.append(text)

        text = "\n".join(slide_text).strip()

        if not text:
            continue

        chunks = chunk_text(
            text=text,
            page_number=slide_number,
            start_index=chunk_index,
        )

        text_chunks.extend(chunks)
        chunk_index += len(chunks)

    return text_chunks, False